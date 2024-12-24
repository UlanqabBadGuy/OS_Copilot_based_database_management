import os
from pptx import Presentation
from pptx.util import Inches
from PIL import Image


def create_ppt_from_folder(folder_path, output_ppt_path):
    """
    Reads all files in a folder and creates a PowerPoint presentation with them.

    Args:
    - folder_path (str): Path to the folder containing files.
    - output_ppt_path (str): Path to save the generated PowerPoint presentation.

    Returns:
    - str: Path of the saved PowerPoint file.
    """
    # Create a presentation
    presentation = Presentation()

    # Iterate over all files in the folder
    for file_name in sorted(os.listdir(folder_path)):
        file_path = os.path.join(folder_path, file_name)
        if os.path.isfile(file_path):
            # Add a slide for each file
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide layout

            # Process text files
            if file_name.lower().endswith(".txt"):
                with open(file_path, "r", encoding="utf-8") as file:
                    content = file.read()
                text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5.5))
                text_frame = text_box.text_frame
                text_frame.text = content

            # Process image files
            elif file_name.lower().endswith((".png", ".jpg", ".jpeg", ".bmp", ".gif")):
                try:
                    img = Image.open(file_path)
                    img_width, img_height = img.size
                    slide_width = presentation.slide_width
                    slide_height = presentation.slide_height

                    # Resize image to fit the slide
                    ratio = min(slide_width / img_width, slide_height / img_height)
                    img_width, img_height = int(img_width * ratio), int(img_height * ratio)
                    left = (slide_width - img_width) / 2
                    top = (slide_height - img_height) / 2

                    slide.shapes.add_picture(file_path, left, top, width=img_width, height=img_height)
                except Exception as e:
                    print(f"Error processing image {file_name}: {e}")

            # Add unsupported file types as text
            else:
                text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5.5))
                text_frame = text_box.text_frame
                text_frame.text = f"Unsupported file: {file_name}"

    # Save the presentation
    presentation.save(output_ppt_path)
    return output_ppt_path